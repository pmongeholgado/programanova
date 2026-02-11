# app/pptx_generator_pro.py

from pptx.enum.shapes import PP_PLACEHOLDER
from io import BytesIO
import base64
from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES_CON_IMAGEN = {1, 6, 9, 10}


def _decode_data_url_to_bytes(data_url: str) -> bytes:
    """
    Convierte "data:image/png;base64,AAAA..." -> bytes reales
    """
    if not data_url or "base64," not in data_url:
        raise ValueError("DataURL inválida o vacía")

    b64_part = data_url.split("base64,", 1)[1].strip()
    return base64.b64decode(b64_part)


def _remove_non_title_placeholders(slide):
    """
    Elimina placeholders que NO sean título (evita el marco/recuadro que tapa la imagen)
    """
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            ph_type = shape.placeholder_format.type
            if ph_type not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                shape.element.getparent().remove(shape.element)
        except Exception:
            pass


def crear_pptx_con_imagenes(titulo: str, slides: list, image_dataurls_by_slide: dict) -> bytes:
    """
    Retorna bytes del PPTX listo para descargar.
    - slides: lista de {"title":..., "bullets":[...]}
    - image_dataurls_by_slide: {1: dataUrl, 6: dataUrl, 9: dataUrl, 10: dataUrl}
    """
    prs = Presentation()

    # Normalizar claves por si vienen como "1","6","9","10"
    if image_dataurls_by_slide:
        image_dataurls_by_slide = {int(k): v for k, v in image_dataurls_by_slide.items()}
    else:
        image_dataurls_by_slide = {}

    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_title_only = prs.slide_layouts[5]

    # --- Slide 1 (portada)
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = titulo or "Presentación Nova"
    if slide1.placeholders and len(slide1.placeholders) > 1:
        try:
            slide1.placeholders[1].text = "Generada con IA - Programa Nova"
        except Exception:
            pass

    # Imagen Slide 1
    if 1 in image_dataurls_by_slide:
        img_bytes = _decode_data_url_to_bytes(image_dataurls_by_slide[1])
        slide1.shapes.add_picture(BytesIO(img_bytes), Inches(1.0), Inches(2.0), width=Inches(8.0))

    # --- Resto de slides
    for idx, s in enumerate(slides, start=1):
        if idx == 1:
            continue

        # Si la slide lleva imagen: Title Only + limpiar placeholders no-título
        if idx in SLIDES_CON_IMAGEN:
            slide = prs.slides.add_slide(layout_title_only)
            _remove_non_title_placeholders(slide)
        else:
            slide = prs.slides.add_slide(layout_content)

        slide.shapes.title.text = s.get("title", f"Diapositiva {idx}")

        # Bullets SOLO en slides sin imagen
        if idx not in SLIDES_CON_IMAGEN:
            body = slide.placeholders[1].text_frame
            body.clear()

            bullets = s.get("bullets", []) or []
            for j, b in enumerate(bullets):
                p = body.paragraphs[0] if j == 0 else body.add_paragraph()
                p.text = str(b)
                p.font.size = Pt(18)
        else:
            # (Opcional) si quieres bullets también en las de imagen, aquí habría que crear textbox.
            pass

        # Insertar imagen SOLO en 6, 9, 10 (y 1 ya se hizo arriba)
        if idx in SLIDES_CON_IMAGEN and idx in image_dataurls_by_slide and idx != 1:
            img_bytes = _decode_data_url_to_bytes(image_dataurls_by_slide[idx])
            slide.shapes.add_picture(BytesIO(img_bytes), Inches(5.5), Inches(2.0), width=Inches(3.8))

    # Guardar PPTX en memoria
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
