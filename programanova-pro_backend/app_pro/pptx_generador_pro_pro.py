# app/pptx_generator_pro.py

from app_pro.ia_assets_pro import generar_imagenes_ppt
from io import BytesIO
import base64
from pptx import Presentation
from pptx.util import Inches, Pt

SLIDES_CON_IMAGEN = {6, 9, 10}

def _decode_data_url_to_bytes(data_url: str) -> bytes:
    """
    Convierte "data:image/png;base64,AAAA..." -> bytes reales
    """
    if not data_url or "base64," not in data_url:
        raise ValueError("DataURL inválida o vacía")

    b64_part = data_url.split("base64,", 1)[1].strip()
    return base64.b64decode(b64_part)

def crear_pptx_con_imagenes(titulo: str, slides: list, image_dataurls_by_slide: dict) -> bytes:
    """
    Retorna bytes del PPTX listo para descargar.
    - slides: lista de {"title":..., "bullets":[...]}
    - image_dataurls_by_slide: {1: dataUrl, 6: dataUrl, 9: dataUrl, 10: dataUrl}
    """
    prs = Presentation()
    
    # Normalizar claves del dict (a veces llegan como "1","6","9","10" desde JSON)
    if image_dataurls_by_slide:
        image_dataurls_by_slide = {int(k): v for k, v in image_dataurls_by_slide.items()}
    else:
        image_dataurls_by_slide = {}

        # BLINDAJE: asegurar que el dict existe siempre y evitar errores si viene None
        image_dataurls_by_slide = image_dataurls_by_slide or {}

    # Layouts típicos: 0=title, 1=title+content, 5=title only
    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_title_only = prs.slide_layouts[5]

    # --- Slide 1 (portada)
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = titulo or "Presentación Nova"
    if slide1.placeholders and len(slide1.placeholders) > 1:
        try:
            slide1.placeholders[1].text = "Generada con IA - Programa Nova"
        except:
            pass

    # Imagen Slide 1
    if 1 in image_dataurls_by_slide:
        img_bytes = _decode_data_url_to_bytes(image_dataurls_by_slide[1])
        slide1.shapes.add_picture(BytesIO(img_bytes), Inches(1.0), Inches(2.0), width=Inches(8.0))

    # --- Resto de slides
    for idx, s in enumerate(slides, start=1):
        # ya usamos slide 1 arriba
        if idx == 1:
            continue

        # layout para texto
        slide = prs.slides.add_slide(layout_content)
        slide.shapes.title.text = s.get("title", f"Diapositiva {idx}")

        body = slide.placeholders[1].text_frame
        body.clear()

        bullets = s.get("bullets", []) or []
        for j, b in enumerate(bullets):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            p.text = str(b)
            p.font.size = Pt(18)

        # Insertar imagen SOLO en 6, 9, 10
        if idx in SLIDES_CON_IMAGEN and idx in image_dataurls_by_slide:
            img_bytes = _decode_data_url_to_bytes(image_dataurls_by_slide[idx])
            slide.shapes.add_picture(BytesIO(img_bytes), Inches(5.5), Inches(2.0), width=Inches(3.8))

    # Guardar PPTX en memoria
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
